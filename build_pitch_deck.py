"""Build SrotaAI judge-pitch deck → deck/SrotaAI_Pitch.pptx.

Run:
    .venv/bin/python build_pitch_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- palette ----------
INK     = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
PAPER   = RGBColor(0xF8, 0xFA, 0xFC)   # slate-50
MUTED   = RGBColor(0x64, 0x74, 0x8B)   # slate-500
SUBTLE  = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
ACCENT  = RGBColor(0x25, 0x63, 0xEB)   # blue-600
EMERALD = RGBColor(0x05, 0x96, 0x69)   # emerald-600
RED     = RGBColor(0xDC, 0x26, 0x26)   # red-600
AMBER   = RGBColor(0xD9, 0x77, 0x06)   # amber-600
VIOLET  = RGBColor(0x7C, 0x3A, 0xED)   # violet-600
INDIGO  = RGBColor(0x4F, 0x46, 0xE5)   # indigo-600
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9 widescreen
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color=PAPER):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill(s, color)
    return s


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=18, color=INK, bullet_color=ACCENT, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        p.space_after = Pt(gap)
        # bullet character
        b = p.add_run()
        b.text = "•  "
        b.font.name = "Calibri"
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        # text
        t = p.add_run()
        t.text = item
        t.font.name = "Calibri"
        t.font.size = Pt(size)
        t.font.color.rgb = color
    return tb


def add_section_label(slide, text):
    # small uppercase eyebrow at top
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.35),
             text.upper(), size=10, bold=True, color=ACCENT)


def add_footer(slide, page, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             "SrotaAI · Real-Time Pharmacovigilance for India",
             size=9, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.3),
             f"{page} / {total}",
             size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_icon_tile(slide, x, y, color, label, value):
    """Pipeline-style coloured rounded square."""
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.2), Inches(1.2))
    tile.adjustments[0] = 0.22
    fill(tile, color)
    add_text(slide, x - Inches(0.4), y + Inches(1.3), Inches(2.0), Inches(0.3),
             label.upper(), size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x - Inches(0.4), y + Inches(1.55), Inches(2.0), Inches(0.5),
             value, size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)


# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 10  # slide count for footer

    # ------------------------------------------------------------
    # SLIDE 1 — Cover / Founder Moment
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)

    # subtle accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(0.5), Inches(0.06))
    fill(bar, ACCENT)

    add_text(s, Inches(0.6), Inches(0.75), Inches(6), Inches(0.4),
             "SROTAAI  ·  AI FOR BHARAT 2026  ·  THEME 6",
             size=10, bold=True, color=SUBTLE)

    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.2),
             "Healthcare signals already exist online.",
             size=44, bold=True, color=WHITE, line_spacing=1.05)

    add_text(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.8),
             "Current systems are not designed to continuously listen to them.",
             size=22, color=SUBTLE)

    add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(1.2),
             "SrotaAI",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             "Real-time pharmacovigilance signal intelligence for India",
             size=14, color=SUBTLE)

    # ------------------------------------------------------------
    # SLIDE 2 — The Founder Story
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "The Origin")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "A product analyst's frustration",
             size=32, bold=True, color=INK)

    add_bullets(s, Inches(0.6), Inches(2.1), Inches(12), Inches(4.2), [
        "As a Product Analyst at a healthcare company, market research is a major part of my work.",
        "While researching public feedback for a healthcare monitoring product, I found myself manually scrolling through Reddit, forums, Google reviews, Facebook groups, and AI search tools — just to understand recurring patient concerns.",
        "The information already existed online.",
        "But there was no intelligent system capable of continuously listening, organizing, validating, and investigating these fragmented healthcare conversations in real time.",
    ], size=18)

    # callout
    cb = add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.7), RGBColor(0xEF, 0xF6, 0xFF))
    cb.line.color.rgb = ACCENT
    cb.line.width = Pt(0.75)
    add_text(s, Inches(0.85), Inches(6.1), Inches(11.7), Inches(0.5),
             "That experience became the starting point for SrotaAI.",
             size=16, bold=True, color=ACCENT)
    add_footer(s, 2, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 3 — Why Existing Systems Fail
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "The Gap")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "Why existing systems fail",
             size=32, bold=True, color=INK)

    # left column — "they do this"
    add_text(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(0.4),
             "TODAY'S WORKFLOWS", size=11, bold=True, color=MUTED)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(5.8), Inches(3.5), [
        "Manual scrolling across dozens of channels",
        "Fragmented across Reddit, forums, reviews, social",
        "Reactive — react to news, not to early chatter",
        "Hard to scale across drugs, geographies, languages",
    ], size=15, bullet_color=RED)

    # right column — "they don't do this"
    add_text(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(0.4),
             "MOST TOOLS STOP AT", size=11, bold=True, color=MUTED)
    add_bullets(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(3.5), [
        "Tracking mentions",
        "Monitoring keywords",
        "Generating dashboards",
    ], size=15, bullet_color=AMBER)

    add_text(s, Inches(7.0), Inches(4.4), Inches(5.8), Inches(0.4),
             "VERY FEW", size=11, bold=True, color=MUTED)
    add_bullets(s, Inches(7.0), Inches(4.8), Inches(5.8), Inches(2), [
        "Validate signals statistically",
        "Connect evidence to claims",
        "Provide explainability",
        "Support investigation workflows",
    ], size=15, bullet_color=EMERALD)

    # closing line
    cb = add_rect(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.55), INK)
    add_text(s, Inches(0.85), Inches(6.48), Inches(11.7), Inches(0.5),
             "The challenge is not lack of data. It's the lack of intelligent listening infrastructure.",
             size=14, bold=True, color=WHITE)
    add_footer(s, 3, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 4 — What is SrotaAI
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Introducing SrotaAI")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "An explainable pharmacovigilance intelligence platform",
             size=28, bold=True, color=INK)

    add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.6),
             "Transforms patient-generated internet data into statistically validated healthcare safety signals.",
             size=15, color=MUTED)

    # 5 pillars in a row of icon tiles
    pillars = [
        ("Monitor",   "Public sources, 24×7",          ACCENT),
        ("Extract",   "Entities & sentiment",          INDIGO),
        ("Correlate", "Drug ↔ symptom pairs",          VIOLET),
        ("Validate",  "PRR / χ² disproportionality",   AMBER),
        ("Investigate","Evidence-backed audit trail",   RED),
    ]
    n = len(pillars)
    col_w = (SLIDE_W - Inches(1.2)) / n
    for i, (title, sub, color) in enumerate(pillars):
        x = Inches(0.6) + col_w * i + (col_w - Inches(2.2)) / 2
        # icon tile
        tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.5), Inches(3.1),
                                  Inches(1.2), Inches(1.2))
        tile.adjustments[0] = 0.22
        fill(tile, color)
        # number badge
        add_text(s, x + Inches(0.5), Inches(3.45), Inches(1.2), Inches(0.6),
                 str(i + 1), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        add_text(s, x, Inches(4.5), Inches(2.2), Inches(0.4),
                 title, size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
        # subtitle
        add_text(s, x, Inches(4.9), Inches(2.2), Inches(0.6),
                 sub, size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # bottom strap
    cb = add_rect(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6), RGBColor(0xEF, 0xF6, 0xFF))
    add_text(s, Inches(0.85), Inches(6.4), Inches(11.7), Inches(0.5),
             "MHRA-aligned thresholds  ·  Tamper-evident hash-chained audit  ·  Open architecture",
             size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, 4, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 5 — Product walkthrough (current state numbers)
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Live Product")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "What's running today",
             size=32, bold=True, color=INK)
    add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
             "Active data pipeline — real numbers from the deployed system.",
             size=14, color=MUTED)

    # pipeline tiles
    stages = [
        ("Sources",  "4",        ACCENT),
        ("Records",  "862,852",  INDIGO),
        ("Entities", "13,370",   VIOLET),
        ("Signals",  "12 strong\n+ 360 exploratory", RED),
    ]
    n = len(stages)
    total_w = Inches(11.0)
    start_x = (SLIDE_W - total_w) / 2
    col_w = total_w / n
    for i, (label, value, color) in enumerate(stages):
        cx = start_x + col_w * i + (col_w - Inches(1.2)) / 2
        tile = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(3.0),
                                  Inches(1.2), Inches(1.2))
        tile.adjustments[0] = 0.22
        fill(tile, color)
        add_text(s, cx - Inches(0.6), Inches(4.4), Inches(2.4), Inches(0.4),
                 label.upper(), size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
        add_text(s, cx - Inches(0.6), Inches(4.8), Inches(2.4), Inches(0.9),
                 value, size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
        # arrow between
        if i < n - 1:
            ax = cx + Inches(1.2) + Inches(0.1)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax,
                                       Inches(3.5), col_w - Inches(1.5), Inches(0.2))
            fill(arrow, SUBTLE)

    add_text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.5),
             "5 connector types live: openFDA · Reddit · RSS · HTML stealth · WhatsApp/X fixtures",
             size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, 5, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 6 — The user workflow (story)
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Workflow")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "From configuration to confirmed signal in 4 steps",
             size=28, bold=True, color=INK)

    steps = [
        ("01", "Create monitoring project",
         "Configure drugs, symptoms, sources, and cadence — no code required."),
        ("02", "Ingest from live sources",
         "SrotaAI continuously pulls public healthcare discussions across channels."),
        ("03", "Signal emerges",
         "PRR, χ², and IC disproportionality with MHRA thresholds promote pairs to signals."),
        ("04", "Investigation workspace",
         "Every signal traces back to evidence, sentiment trends, and source attribution."),
    ]
    y = Inches(2.1)
    row_h = Inches(1.05)
    for i, (num, title, sub) in enumerate(steps):
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + i * row_h,
                                   Inches(0.8), Inches(0.8))
        fill(badge, ACCENT)
        add_text(s, Inches(0.6), y + i * row_h + Inches(0.18),
                 Inches(0.8), Inches(0.5),
                 num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title + sub
        add_text(s, Inches(1.7), y + i * row_h - Inches(0.05),
                 Inches(11), Inches(0.4),
                 title, size=18, bold=True, color=INK)
        add_text(s, Inches(1.7), y + i * row_h + Inches(0.4),
                 Inches(11), Inches(0.5),
                 sub, size=13, color=MUTED)
    add_footer(s, 6, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 7 — Why SrotaAI is different (positioning)
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Positioning")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "Signal intelligence — not just monitoring",
             size=28, bold=True, color=INK)

    # Two-column comparison table
    col_x_left  = Inches(0.6)
    col_x_right = Inches(7.0)
    col_w_each  = Inches(5.7)
    head_h      = Inches(0.55)

    # Headers
    h1 = add_rect(s, col_x_left,  Inches(2.0), col_w_each, head_h, RGBColor(0xF1, 0xF5, 0xF9))
    add_text(s, col_x_left + Inches(0.2), Inches(2.1), col_w_each, head_h,
             "Traditional Monitoring", size=13, bold=True, color=MUTED)
    h2 = add_rect(s, col_x_right, Inches(2.0), col_w_each, head_h, ACCENT)
    add_text(s, col_x_right + Inches(0.2), Inches(2.1), col_w_each, head_h,
             "SrotaAI", size=13, bold=True, color=WHITE)

    rows = [
        ("Keyword tracking",      "Signal intelligence"),
        ("Dashboards",            "Investigation workflows"),
        ("Static reporting",      "Continuous monitoring"),
        ("Isolated mentions",     "Explainable correlations"),
        ("Black-box alerts",      "Evidence-backed reasoning"),
    ]
    row_h = Inches(0.65)
    for i, (l, r) in enumerate(rows):
        ry = Inches(2.55) + row_h * i
        bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xF8, 0xFA, 0xFC)
        add_rect(s, col_x_left,  ry, col_w_each, row_h, bg)
        add_rect(s, col_x_right, ry, col_w_each, row_h, RGBColor(0xEF, 0xF6, 0xFF))
        add_text(s, col_x_left + Inches(0.25), ry + Inches(0.1),
                 col_w_each, row_h,
                 l, size=14, color=INK)
        add_text(s, col_x_right + Inches(0.25), ry + Inches(0.1),
                 col_w_each, row_h,
                 r, size=14, bold=True, color=ACCENT)
    add_footer(s, 7, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 8 — Where we are honest (limitations + roadmap)
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Honest Roadmap")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "What's a prototype today, what's production tomorrow",
             size=26, bold=True, color=INK)

    # left: today's limits
    add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
             "TODAY (PROTOTYPE LIMITATIONS)", size=11, bold=True, color=AMBER)
    add_bullets(s, Inches(0.6), Inches(2.4), Inches(6), Inches(4.2), [
        "Symptom dictionary is curated, not full MedDRA — small entity coverage limits record-to-signal yield.",
        "Anonymous Reddit & openFDA only — no Twitter/X firehose, no WhatsApp Business at scale.",
        "Signal trends require a 2nd run to compute (rising / falling) — first run shows 'preliminary'.",
        "Hindi / regional language NER not yet wired (English + transliteration only).",
        "Adding a new connector type today still needs a Python class.",
    ], size=14, bullet_color=AMBER)

    # right: production roadmap
    add_text(s, Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
             "PRODUCTION ROADMAP", size=11, bold=True, color=EMERALD)
    add_bullets(s, Inches(7.0), Inches(2.4), Inches(6), Inches(4.2), [
        "Full MedDRA SOC/PT integration once we receive the official MSSO key — 26K terms, real-world recall.",
        "Agentic source onboarder: an LLM writes new BaseConnector subclasses on demand from a URL + sample.",
        "Multilingual NER (Hindi / Tamil / Bengali / Marathi) via IndicBERT-style fine-tunes.",
        "Graph-based signal intelligence — surface indirect drug-class effects across 1Cr+ records.",
        "Open-by-default monitoring: track a drug with no symptom hypothesis (anomaly-driven).",
    ], size=14, bullet_color=EMERALD)
    add_footer(s, 8, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 9 — Why this matters
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_section_label(s, "Impact")
    add_text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.9),
             "Why earlier signal visibility matters",
             size=28, bold=True, color=INK)

    cards = [
        ("Faster investigations",
         "PV teams stop firefighting; they investigate before signals become news."),
        ("Better workflows",
         "Triage, evidence, and audit in one place — replaces spreadsheets and email chains."),
        ("Proactive monitoring",
         "Continuous listening across India's fragmented healthcare conversation surface."),
        ("Patient intelligence",
         "Explainable signals that regulators, hospitals, and pharma can defensibly act on."),
    ]
    grid_x = [Inches(0.6), Inches(7.0), Inches(0.6), Inches(7.0)]
    grid_y = [Inches(2.0), Inches(2.0), Inches(4.6), Inches(4.6)]
    card_w, card_h = Inches(5.7), Inches(2.3)
    for (title, body), x, y in zip(cards, grid_x, grid_y):
        card = add_rect(s, x, y, card_w, card_h, WHITE)
        card.line.color.rgb = SUBTLE
        card.line.width = Pt(0.75)
        # accent bar
        bar = add_rect(s, x, y, Inches(0.08), card_h, ACCENT)
        add_text(s, x + Inches(0.4), y + Inches(0.3), card_w - Inches(0.6), Inches(0.5),
                 title, size=18, bold=True, color=INK)
        add_text(s, x + Inches(0.4), y + Inches(0.95), card_w - Inches(0.6), card_h - Inches(1.0),
                 body, size=13, color=MUTED)
    add_footer(s, 9, TOTAL)

    # ------------------------------------------------------------
    # SLIDE 10 — Closing
    # ------------------------------------------------------------
    s = add_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, INK)
    add_text(s, Inches(0.6), Inches(0.6), Inches(6), Inches(0.4),
             "ONE LINE",
             size=10, bold=True, color=SUBTLE)

    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(2.4),
             "SrotaAI transforms fragmented patient conversations\ninto explainable healthcare intelligence.",
             size=36, bold=True, color=WHITE, line_spacing=1.15)

    add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
             "github.com/KritikaTri/srota-ai", size=14, color=ACCENT)
    add_text(s, Inches(0.6), Inches(5.85), Inches(12), Inches(0.5),
             "AI for Bharat 2026  ·  Theme 6  ·  Real-Time Social Listening for Patient Safety",
             size=12, color=SUBTLE)

    # ------------------------------------------------------------
    out_dir = Path(__file__).parent / "deck"
    out_dir.mkdir(exist_ok=True)
    out = out_dir / "SrotaAI_Pitch.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
